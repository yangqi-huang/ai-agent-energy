import hashlib
import re
import shutil
import subprocess
import tempfile
from concurrent.futures import ThreadPoolExecutor
from io import BytesIO
from pathlib import Path
from zipfile import ZipFile

from PIL import Image, ImageEnhance, ImageOps
from PyPDF2 import PdfReader, PdfWriter
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from config import (
    OCR_MAX_IMAGES_PER_PAGE,
    OCR_MAX_IMAGES_PER_FILE,
    OCR_MAX_RENDERED_PDF_PAGES,
    OCR_MIN_IMAGE_PIXELS,
    OCR_RENDERED_PDF_PAGE_SIZE,
    OCR_SPARSE_PAGE_CHARS,
    OCR_WORKERS,
)
from material_metrics import extract_key_information_context


def _is_substantial_image(image_bytes: bytes) -> bool:
    try:
        with Image.open(BytesIO(image_bytes)) as image:
            width, height = image.size
        return width * height >= OCR_MIN_IMAGE_PIXELS and min(width, height) >= 180
    except Exception:
        return False


def _preprocess_image(image_bytes: bytes) -> BytesIO:
    image = Image.open(BytesIO(image_bytes))
    image = ImageOps.exif_transpose(image).convert("L")
    width, height = image.size

    if max(width, height) < 1800:
        scale = min(3, 1800 / max(width, height))
        image = image.resize(
            (int(width * scale), int(height * scale)),
            Image.Resampling.LANCZOS,
        )

    image = ImageEnhance.Contrast(image).enhance(1.8)
    image = ImageOps.autocontrast(image)

    output = BytesIO()
    image.save(output, format="PNG")
    output.seek(0)
    return output


def _ocr_quality_score(text: str) -> int:
    if not text:
        return 0
    useful_chars = re.findall(r"[\w\u4e00-\u9fff%$€£¥./,-]", text)
    digits = re.findall(r"\d", text)
    key_terms = re.findall(
        r"IRR|NPV|CAPEX|OPEX|EBITDA|Reserve|Area|储量|面积|投资|回收期|产能|万吨|合作",
        text,
        re.IGNORECASE,
    )
    replacement_noise = text.count("�") + text.count("?")
    return len(useful_chars) + len(digits) * 2 + len(key_terms) * 12 - replacement_noise * 5


def _best_ocr_result(results: list[str]) -> str:
    return max(results, key=_ocr_quality_score, default="")


LOW_VALUE_TEXT_PATTERNS = [
    r"strictly\s+confiden\s*tial",
    r"may\s+not\s+be\s+distribut",
    r"prior\s+writt\s*en\s+consen\s*t",
    r"this\s+documen\s*t",
]


def _is_low_value_pdf_text(page_text: str) -> bool:
    normalized = re.sub(r"\s+", " ", page_text).strip().lower()
    if not normalized:
        return True

    matches = sum(
        1
        for pattern in LOW_VALUE_TEXT_PATTERNS
        if re.search(pattern, normalized, re.IGNORECASE)
    )
    words = re.findall(r"[a-zA-Z\u4e00-\u9fff]{3,}", normalized)
    numbers = re.findall(r"\d", normalized)
    useful_terms = re.findall(
        r"palas|prospek|blok|migas|minarak|brantas|cadangan|sumber|daya|luas|"
        r"reserv|resource|area|prospect|exploration|gas|oil|irr|npv|invest|"
        r"produksi|production|psc|working\s+interest|operator",
        normalized,
        re.IGNORECASE,
    )

    if matches >= 2 and len(useful_terms) <= 2 and len(numbers) <= 2:
        return True
    if len(words) <= 12 and not useful_terms:
        return True
    return False


def ocr_image(image_bytes: bytes, lang: str = "chi_sim+eng") -> tuple[str, str | None]:
    try:
        processed = _preprocess_image(image_bytes)
    except Exception as exc:
        return "", f"图片预处理失败：{exc}"

    try:
        import pytesseract

        results = []
        for config in ["--psm 6", "--psm 11"]:
            text = pytesseract.image_to_string(
                Image.open(processed),
                lang=lang,
                config=config,
                timeout=45,
            )
            if text.strip():
                results.append(text.strip())
        return _best_ocr_result(results).strip(), None
    except Exception as python_error:
        command = shutil.which("tesseract")
        if not command:
            return "", "未检测到OCR引擎，图片文字未提取。"

        with tempfile.NamedTemporaryFile(suffix=".png") as image_file:
            image_file.write(processed.getvalue())
            image_file.flush()
            try:
                result = subprocess.run(
                    [command, image_file.name, "stdout", "-l", lang, "--psm", "6"],
                    check=True,
                    capture_output=True,
                    text=True,
                    timeout=45,
                )
                return result.stdout.strip(), None
            except Exception as cli_error:
                return "", f"OCR失败：{python_error}；{cli_error}"


def _ocr_unique_images(
    images: list[bytes],
    max_images: int = OCR_MAX_IMAGES_PER_PAGE,
) -> list[tuple[str, str | None]]:
    selected = []
    seen = set()

    for image_bytes in images:
        digest = hashlib.sha1(image_bytes).hexdigest()
        if digest in seen or not _is_substantial_image(image_bytes):
            continue
        seen.add(digest)
        selected.append(image_bytes)
        if len(selected) >= max_images:
            break

    if not selected:
        return []

    with ThreadPoolExecutor(max_workers=OCR_WORKERS) as executor:
        return list(executor.map(ocr_image, selected))


def _render_pdf_page_with_pymupdf(file_path: str, page_index: int) -> bytes | None:
    try:
        import fitz
    except Exception:
        return None

    try:
        document = fitz.open(file_path)
        page = document.load_page(page_index)
        pixmap = page.get_pixmap(matrix=fitz.Matrix(2.4, 2.4), alpha=False)
        image = pixmap.tobytes("png")
        document.close()
        return image
    except Exception:
        return None


def _render_pdf_page_with_quicklook(
    reader: PdfReader,
    page_index: int,
    temp_dir: str,
) -> bytes | None:
    command = shutil.which("qlmanage")
    if not command:
        return None

    page_pdf = Path(temp_dir) / f"page_{page_index + 1}.pdf"
    output_dir = Path(temp_dir) / "preview"
    output_dir.mkdir(exist_ok=True)

    writer = PdfWriter()
    writer.add_page(reader.pages[page_index])
    with page_pdf.open("wb") as handle:
        writer.write(handle)

    try:
        subprocess.run(
            [
                command,
                "-t",
                "-s",
                str(OCR_RENDERED_PDF_PAGE_SIZE),
                "-o",
                str(output_dir),
                str(page_pdf),
            ],
            check=True,
            capture_output=True,
            text=True,
            timeout=30,
        )
    except Exception:
        return None

    candidates = sorted(output_dir.glob(f"{page_pdf.name}*.png"))
    if not candidates:
        candidates = sorted(
            output_dir.glob("*.png"),
            key=lambda item: item.stat().st_mtime,
        )
    if not candidates:
        return None
    return candidates[-1].read_bytes()


def _render_pdf_page_to_image(
    file_path: str,
    reader: PdfReader,
    page_index: int,
    temp_dir: str,
) -> bytes | None:
    rendered = _render_pdf_page_with_pymupdf(file_path, page_index)
    if rendered:
        return rendered
    return _render_pdf_page_with_quicklook(reader, page_index, temp_dir)


def _parse_pdf(file_path: str) -> str:
    text = []
    ocr_image_count = 0
    rendered_page_count = 0
    reader = PdfReader(file_path)
    with tempfile.TemporaryDirectory() as temp_dir:
        for index, page in enumerate(reader.pages, start=1):
            page_text = (page.extract_text() or "").strip()
            page_parts = [page_text] if page_text else []
            needs_ocr = (
                len(page_text) < OCR_SPARSE_PAGE_CHARS
                or _is_low_value_pdf_text(page_text)
            )

            should_render_page = (
                needs_ocr
                and rendered_page_count < OCR_MAX_RENDERED_PDF_PAGES
                and ocr_image_count < OCR_MAX_IMAGES_PER_FILE
            )
            rendered_ocr_ok = False
            if should_render_page:
                rendered = _render_pdf_page_to_image(file_path, reader, index - 1, temp_dir)
                if rendered:
                    rendered_page_count += 1
                    ocr_image_count += 1
                    recognized, error = ocr_image(rendered)
                    if recognized:
                        rendered_ocr_ok = True
                        page_parts.append(f"【整页渲染OCR】\n{recognized}")
                    elif error:
                        page_parts.append(f"【整页渲染OCR】{error}")

            if (
                needs_ocr
                and not rendered_ocr_ok
                and ocr_image_count < OCR_MAX_IMAGES_PER_FILE
            ):
                try:
                    image_bytes = [image.data for image in page.images]
                except Exception:
                    image_bytes = []
                results = _ocr_unique_images(
                    image_bytes,
                    min(
                        OCR_MAX_IMAGES_PER_PAGE,
                        OCR_MAX_IMAGES_PER_FILE - ocr_image_count,
                    ),
                )
                ocr_image_count += len(results)
                for image_index, (recognized, error) in enumerate(results, start=1):
                    if recognized:
                        page_parts.append(f"【页面图片{image_index} OCR】\n{recognized}")
                    elif error:
                        page_parts.append(f"【页面图片{image_index}】{error}")

            if needs_ocr and len("\n".join(page_parts).strip()) < OCR_SPARSE_PAGE_CHARS:
                page_parts.append("【OCR提示】本页原生文本较少，整页渲染或图片OCR未提取到足够文字。")

            if page_parts:
                text.append(f"【PDF第{index}页】\n" + "\n\n".join(page_parts))

        if rendered_page_count >= OCR_MAX_RENDERED_PDF_PAGES:
            text.append(
                f"【提示】PDF整页OCR已达到上限 {OCR_MAX_RENDERED_PDF_PAGES} 页，"
                "后续页面保留可提取文本。"
            )
    return "\n\n".join(text)


def _parse_docx(file_path: str) -> str:
    document = Document(file_path)
    text = [paragraph.text for paragraph in document.paragraphs if paragraph.text]

    for index, table in enumerate(document.tables, start=1):
        rows = []
        for row in table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            text.append(f"【Word表格 {index}】\n" + "\n".join(rows))

    try:
        with ZipFile(file_path) as archive:
            images = [
                archive.read(name)
                for name in archive.namelist()
                if name.startswith("word/media/")
            ]
        for image_index, (recognized, error) in enumerate(
            _ocr_unique_images(images),
            start=1,
        ):
            if recognized:
                text.append(f"【Word图片{image_index} OCR】\n{recognized}")
            elif error:
                text.append(f"【Word图片{image_index}】{error}")
    except Exception:
        pass

    return "\n\n".join(text)


def _parse_pptx(file_path: str) -> str:
    presentation = Presentation(file_path)
    output = []
    ocr_image_count = 0

    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_parts = []
        slide_images = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_parts.append(shape.text.strip())

            if getattr(shape, "has_table", False):
                rows = []
                for row in shape.table.rows:
                    cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    if any(cells):
                        rows.append(" | ".join(cells))
                if rows:
                    slide_parts.append("【表格】\n" + "\n".join(rows))

            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                slide_images.append(shape.image.blob)

        results = []
        if ocr_image_count < OCR_MAX_IMAGES_PER_FILE:
            results = _ocr_unique_images(
                slide_images,
                min(
                    OCR_MAX_IMAGES_PER_PAGE,
                    OCR_MAX_IMAGES_PER_FILE - ocr_image_count,
                ),
            )
            ocr_image_count += len(results)

        for image_index, (recognized, error) in enumerate(results, start=1):
            if recognized:
                slide_parts.append(f"【图片{image_index} OCR】\n{recognized}")
            elif error:
                slide_parts.append(f"【图片{image_index}】{error}")

        try:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_parts.append(f"【备注】\n{notes}")
        except Exception:
            pass

        if slide_parts:
            output.append(f"【PPT第{slide_index}页】\n" + "\n\n".join(slide_parts))

    return "\n\n".join(output)


def _parse_image(file_path: str, file_name: str) -> str:
    text, error = ocr_image(Path(file_path).read_bytes())
    if text:
        return f"【图片文件：{file_name}】\n{text}"
    return f"【图片文件：{file_name}】\n{error or '未识别到图片文字'}"


def parse_file(file_path: str, file_name: str) -> str:
    extension = Path(file_name).suffix.lower()

    try:
        if extension == ".pdf":
            return _parse_pdf(file_path)
        if extension == ".docx":
            return _parse_docx(file_path)
        if extension == ".pptx":
            return _parse_pptx(file_path)
        if extension in {".png", ".jpg", ".jpeg"}:
            return _parse_image(file_path, file_name)
        return f"不支持的文件类型：{extension}"
    except Exception as exc:
        return f"解析失败：{exc}"


def parse_uploaded_files(files: list[dict]) -> str:
    blocks = []
    for file_info in files:
        name = file_info["name"]
        content = parse_file(file_info["datapath"], name)
        blocks.append(f"=== 文件名：{name} ===\n\n{content}")
    combined = "\n\n" + ("\n\n" + "=" * 50 + "\n\n").join(blocks)
    key_context = extract_key_information_context(combined)
    if key_context:
        return "\n\n".join([key_context, combined])
    return combined
