from email.mime import text

from shiny import App, Inputs, Outputs, Session, reactive, ui, render
import requests
from pathlib import Path

# 文档解析
from docx import Document
from PyPDF2 import PdfReader
from pptx import Presentation

# 搜索
from duckduckgo_search import DDGS

# ======================================
# UI 界面
# ======================================

app_ui = ui.page_sidebar(

    # Sidebar
    ui.sidebar(
        ui.input_password(
            "api_key",
            "DeepSeek API Key",
            placeholder="输入你的 API Key"
        ),

        ui.hr(),

        ui.h5("📁 上传文件"),

        ui.input_file(
            "upload",
            "上传 PPT / PDF / Word / 图片",
            multiple=True,
            accept=[
                ".ppt",
                ".pptx",
                ".pdf",
                ".docx",
                ".png",
                ".jpg",
                ".jpeg"
            ]
        ),

        ui.hr(),

        ui.input_action_button(
            "run_agent",
            "▶️ 开始解析+结构化",
            class_="btn-primary w-100"
        ),

        width=320
    ),

    # 页面主体
    ui.tags.head(
        ui.tags.title("项目简报信息排版AI Agent")
    ),

    ui.card(

        ui.h4("✅ AI 结构化输出"),

        ui.layout_columns(

            # ======================================
            # 左侧：可编辑 Markdown
            # ======================================

            ui.card(

                ui.h5("✏️ 编辑区"),

                ui.input_text_area(
                    "editable_ai_result",
                    label=None,
                    value="",
                    width="100%",
                    height="700px"
                ),

                full_screen=True

            ),

            # ======================================
            # 右侧：Markdown 实时预览
            # ======================================

            ui.card(

                ui.h5("📖 实时预览"),

                ui.div(

                    ui.output_ui("markdown_preview"),

                    style="""
                    height:700px;
                    overflow-y:auto;
                    padding:15px;
                    border:1px solid #ddd;
                    border-radius:8px;
                    background:#fff;
                    """
                ),

                full_screen=True
            ),

            col_widths=[6, 6]

        ),

        height="100%"
    ),

    fillable=True
)

# ======================================
# 服务端逻辑
# ======================================

def server(input: Inputs, output: Outputs, session: Session):

    # Reactive 状态
    extracted_text_value = reactive.Value("")
    ai_output_result_value = reactive.Value("等待AI处理...")

    # ======================================
    # DuckDuckGo 搜索
    # ======================================
    def web_search(query: str, max_results=5):

        text = ""

        try:

            with DDGS() as ddgs:

                results = ddgs.text(
                    query,
                    max_results=max_results
                )

                for r in results:

                    title = r.get("title", "")
                    body = r.get("body", "")
                    body = body[:800]
                    href = r.get("href", "")

                    text += f"""
    标题：
    {title}

    内容：
    {body}

    链接：
    {href}

    -------------------------
    """

        except Exception as e:

            text += f"搜索失败：{str(e)}"

        return text

    # ======================================
    # EIA API：通用数据获取 + 格式化
    # ======================================

    def format_eia_data(result: dict, title: str = "EIA 行业数据", max_rows: int = 10) -> str:
        """
        将不同 EIA API 返回结果统一格式化成适合喂给 LLM 的文本。
        兼容大多数 EIA v2 API 返回结构。
        """

        try:
            response = result.get("response", {})
            records = response.get("data", [])

            if not records:
                return f"{title}\n\n未获取到有效 EIA 数据。\n"

            text = f"# {title}\n\n"
            text += f"数据来源：U.S. Energy Information Administration (EIA)\n\n"

            # 只取前 max_rows 条，避免 token 爆炸
            for idx, row in enumerate(records[:max_rows], start=1):
                text += f"## 数据记录 {idx}\n"

                # 常见核心字段，优先展示
                preferred_fields = [
                    "period",
                    "value",
                    "units",
                    "series-description",
                    "series",
                    "product-name",
                    "process-name",
                    "area-name",
                    "duoarea",
                    "frequency",
                ]

                shown_keys = set()

                for key in preferred_fields:
                    if key in row and row.get(key) not in [None, ""]:
                        text += f"- {key}: {row.get(key)}\n"
                        shown_keys.add(key)

                # 兜底：展示其他字段
                for key, value in row.items():
                    if key not in shown_keys and value not in [None, ""]:
                        text += f"- {key}: {value}\n"

                text += "\n"

            return text

        except Exception as e:
            return f"{title}\n\nEIA 数据格式化失败：{str(e)}\n"


    def get_eia_data(eia_url: str, title: str = "EIA 行业数据") -> str:
        """
        请求 EIA API，并返回格式化后的文本。
        """

        try:
            resp = requests.get(
                eia_url,
                timeout=30
            )

            resp.raise_for_status()

            result = resp.json()

            return format_eia_data(
                result=result,
                title=title,
                max_rows=10
            )

        except Exception as e:
            return f"# {title}\n\nEIA 数据获取失败：{str(e)}\n"

    # ======================================
    # 文件解析
    # ======================================

    def parse_file(file_path: str, file_name: str):

        ext = Path(file_name).suffix.lower()
        text = ""

        try:

            # PDF
            if ext == ".pdf":

                reader = PdfReader(file_path)

                for page in reader.pages:

                    t = page.extract_text()

                    if t:
                        text += t + "\n\n"

            # Word
            elif ext == ".docx":

                doc = Document(file_path)

                for para in doc.paragraphs:
                    text += para.text + "\n"

            # PPT
            elif ext == ".pptx":

                prs = Presentation(file_path)

                for slide in prs.slides:

                    for shape in slide.shapes:

                        if hasattr(shape, "text"):

                            if shape.text:
                                text += shape.text + "\n"

            # 图片
            elif ext in [".png", ".jpg", ".jpeg"]:

                text += f"【图片文件】{file_name}\n"
                text += "当前版本暂未OCR识别\n"

            else:

                text += f"不支持的文件类型：{ext}\n"

        except Exception as e:

            text += f"解析失败：{str(e)}\n"

        return text

    # ======================================
    # 点击按钮执行
    # ======================================

    @reactive.Effect
    @reactive.event(input.run_agent)
    def process_all():

        api_key = input.api_key().strip()
        files = input.upload()

        # 校验
        if not api_key:

            ui.notification_show(
                "请输入 API Key",
                type="error"
            )

            return

        if not files:

            ui.notification_show(
                "请上传文件",
                type="error"
            )

            return

        # 进度条
        with ui.Progress(min=0, max=100) as p:

            p.set(10, message="正在解析文件...")

            all_text = ""

            # 遍历文件
            for f in files:

                path = f["datapath"]
                name = f["name"]

                all_text += f"\n=== 文件名：{name} ===\n\n"

                parsed_text = parse_file(path, name)

                all_text += parsed_text

                all_text += "\n" + "=" * 50 + "\n"

            # 更新提取文本
            extracted_text_value.set(all_text.strip())

            p.set(30, message="正在生成搜索关键词...")

            # ======================================
            # 自动生成搜索关键词
            # ======================================

            keyword_prompt = f"""
            你是专业项目研究员。

            请根据项目材料：

            生成以下类型的联网搜索计划：

            1. 行业市场
            2. 政策法规
            3. 合作方背景
            4. 区域经济
            5. 风险因素

            每类给出2个搜索关键词。
            
            要求：
            - 不要解释
            - 不要编号

            项目材料：
            {all_text}
            """

            keyword_resp = requests.post(

                url="https://api.deepseek.com/chat/completions",

                headers={
                    "Authorization": f"Bearer {api_key}",
                    "Content-Type": "application/json"
                },

                json={

                    "model": "deepseek-v4-flash",

                    "messages": [

                        {
                            "role": "system",
                            "content": "你是专业项目关键词提取助手"
                        },

                        {
                            "role": "user",
                            "content": keyword_prompt
                        }
                    ],

                    "temperature": 0.1
                }
            )

            keyword_result = keyword_resp.json()

            keywords_text = keyword_result["choices"][0]["message"]["content"]

            keywords = keywords_text.split("\n")

            p.set(50, message="正在联网搜索补充信息...")

            search_result_text = ""

            # ======================================
            # EIA 行业数据
            # ======================================

            eia_result = ""

            if (
                "lng" in all_text.lower()
                or "lpg" in all_text.lower()
                or "天然气" in all_text
                or "液化天然气" in all_text
            ):

                p.set(60, message="正在获取 EIA 行业数据...")

                eia_url = "https://api.eia.gov/v2/international/data/?frequency=annual&data[0]=value&start=1949&end=2026&sort[0][column]=period&sort[0][direction]=desc&offset=0&length=5000"

                eia_result = get_eia_data(
                    eia_url=eia_url,
                    title="EIA 天然气 / LNG 行业数据"
                )

            for kw in keywords:

                kw = kw.strip()

                if kw:

                    search_result_text += f"""

            =========================
            搜索关键词：
            {kw}
            =========================

            """

                    result = web_search(kw)

                    search_result_text += result

            # ======================================
            # 拼接最终 Prompt
            # ======================================

            final_user_prompt = f"""

            # 用户上传的项目材料

            {all_text}

            # 联网搜索补充资料

            {search_result_text}

            # EIA 行业数据

            {eia_result}

            # 任务要求

            请结合：

            1. 用户上传材料
            2. 联网搜索资料
            3. EIA 行业数据（如果有）

            输出完整项目简报。

            注意：

            - 优先采用用户原始材料
            - 搜索结果仅作为补充参考
            - EIA 数据属于高可信行业数据
            - 若搜索结果与原文冲突，请明确标注“待核实”
            - 不要编造不存在的数据
            - 请保持专业投研风格
            """

            p.set(70, message="正在调用 DeepSeek...")

            # System Prompt
            system_prompt = """
            你是专业的【广东九丰能源项目简报信息排版AI Agent】。

            请结合你所搜查到的信息，严格按照以下9大模块以及相应要求输出：

            要求：
            一、若有些模块在原文件中没有相应的信息对应，请跳过
            二、请使用Markdown格式输出。
                要求：
                - 一级标题使用 #
                - 二级标题使用 ##
                - 表格使用 markdown table
                - 风险项使用 bullet points
                - 内容适合直接生成项目简报

            9大模块框架：    
            1. 项目摘要
                a. 一句话概括：包括但不限于项目所在国家、项目类型、资源或资产属性、拟建设内容、合作方式等
                b. 核心投资逻辑：投资指标（如capex与收益）、合作方式与商业模式（如股权投资、联合开发等）、上下游（如区位靠近港口/工业园等）、核心待核实项（如资源真实性、权属清晰度等）
            2. 项目定位
                a. 项目所在国家、所在区域、项目属性等
                b. 归属（如上游区块、GPP/LNG类项目、收购类、投标类）
                c. 资源基础（气量分析、气质分析）
                d. 基础设施（基础设施情况、设备情况）
            3. 项目合作方
                a. 基本信息
                b. 主营业务
                c. 核心资产
                d. 财务状况
                e. 近五年重大决策
            4. 市场消纳
                a. 产品销售对象（如管道气——电厂、工业园；LNG——槽车半径；LPG——民用、商业）
                b. 潜在买方
                c. 竞争分析（如周边管道气供应商、LNG进口商等）
            5. 投资方式
                a. 参与逻辑
                b. 投资形式（如建设类、收购类、投标类等）
                c. 商业模式（如盈利来源与风险特征、合作架构等）
                d. 投资指标（如capex与收益、投资回收期、IRR等）
            6. 风险识别
                a. 汇兑风险（如外币支付、外币收入等）
                b. 税务风险（如税率、税收优惠政策、转让定价等）
                c. 市场风险（如市场需求、价格波动、竞争格局等）
                d. 运营风险（如上游资源风险、基础设施风险、合作方风险等）
            7. 结论与建议下一步计划

            最后必须输出：

            【人工审慎把关】判断：
            【思考逻辑】：
            """

            try:

                resp = requests.post(

                    url="https://api.deepseek.com/chat/completions",

                    headers={
                        "Authorization": f"Bearer {api_key}",
                        "Content-Type": "application/json"
                    },

                    json={

                        "model": "deepseek-v4-pro",

                        "messages": [

                            {
                                "role": "system",
                                "content": system_prompt
                            },

                            {
                                "role": "user",
                                "content": final_user_prompt
                            }
                        ],

                        "temperature": 0.2
                    },

                    timeout=120
                )

                resp.raise_for_status()

                result = resp.json()

                ai_output = result["choices"][0]["message"]["content"]

                # 更新 AI 输出
                ai_output_result_value.set(ai_output)

                ui.update_text_area(
                "editable_ai_result",
                value=ai_output
                )

                p.set(100, message="完成")

                ui.notification_show(
                    "✅ 处理完成",
                    type="message"
                )

            except Exception as e:

                ai_output_result_value.set(
                    f"请求失败：{str(e)}"
                )

                ui.notification_show(
                    f"❌ 错误：{str(e)}",
                    type="error"
                )

    # ======================================
    # 输出渲染
    # ======================================

    @output
    @render.ui
    def markdown_preview():

        content = input.editable_ai_result()

        return ui.markdown(content)

    @output
    @render.text
    def extracted_text():

        return extracted_text_value()

    @output
    @render.text
    def ai_result():

        return ai_output_result_value()

# ======================================
# 启动 App
# ======================================

app = App(app_ui, server)